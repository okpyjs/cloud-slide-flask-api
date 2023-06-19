import os
from typing import List

import aspose.slides as slides
from google.oauth2 import service_account
from googleapiclient.discovery import build
from pptx import Presentation

from utils.Drive import Drive


class Slide:
    def __init__(self) -> None:

        SERVICE_ACCOUNT_FILE = "masayuki_service.json"
        SCOPES = ["https://www.googleapis.com/auth/presentations"]
        credentials = service_account.Credentials.from_service_account_file(
            SERVICE_ACCOUNT_FILE, scopes=SCOPES
        )
        self.service = build("slides", "v1", credentials=credentials)

    def convert(self, file: str):
        return

    def merge(self, presentation_ids: list[str], result_file_name: str):
        # IDs of the Google Slides to merge
        for id in presentation_ids:
            Drive().download(id)

        prs = Presentation(f"{presentation_ids[0]}.pptx")
        for id in presentation_ids[1:]:
            presentation = Presentation(f"{id}.pptx")
            for slide in presentation.slides:
                new_slide = prs.slides.add_slide(slide.slide_layout)
                for shape in slide.shapes:
                    new_shape = new_slide.shapes._spTree.append(shape.element)

        prs.save(result_file_name)

        file_info = Drive().upload(result_file_name)

        return file_info

    def merge1(self, presentation_ids: list[str], result_file_name: str):
        # IDs of the Google Slides to merge
        for id in presentation_ids:
            Drive().download(id)

        with slides.Presentation(f"{presentation_ids[0]}.pptx") as pres1:

            # Loop through presentation IDs
            for presentation_id in presentation_ids[1:]:
                # Open each PPT using presentation ID
                with slides.Presentation(f"{presentation_id}.pptx") as pres:
                    # Loop through slides in the current PPT
                    for slide in pres.slides:
                        # Clone slide to the first PPT
                        pres1.slides.add_clone(slide)

            # Save merged PPT
            pres1.save(result_file_name, slides.export.SaveFormat.PPTX)

        file_info = Drive().upload(result_file_name)

        return file_info
